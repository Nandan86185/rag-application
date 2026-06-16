import os
import csv
import json


SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".txt", ".md",
    ".csv", ".xlsx", ".xls",
    ".pptx", ".html", ".htm", ".json"
}


def parse_document(file_path: str) -> str:
    """
    Universal document parser.
    Supports: PDF, DOCX, TXT, MD, CSV, XLSX, PPTX, HTML, JSON
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()
    print(f"Parsing document [{ext}]: {file_path}")

    parsers = {
        ".pdf":  _parse_pdf,
        ".docx": _parse_docx,
        ".txt":  _parse_txt,
        ".md":   _parse_txt,
        ".csv":  _parse_csv,
        ".xlsx": _parse_xlsx,
        ".xls":  _parse_xlsx,
        ".pptx": _parse_pptx,
        ".html": _parse_html,
        ".htm":  _parse_html,
        ".json": _parse_json,
    }

    parser_fn = parsers.get(ext)
    if not parser_fn:
        raise ValueError(
            f"Unsupported file type: '{ext}'. "
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    text = parser_fn(file_path)
    if not text or not text.strip():
        raise ValueError(f"No text could be extracted from '{os.path.basename(file_path)}'.")

    print(f"  Extracted {len(text)} characters.")
    return text


# ─────────────────────────────────────────────────────────────────
# Individual parsers
# ─────────────────────────────────────────────────────────────────

def _parse_pdf(file_path: str) -> str:
    # Attempt 1: pdfplumber
    try:
        import pdfplumber
        pages = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text and text.strip():
                    pages.append(text.strip())
        if pages:
            print(f"  pdfplumber extracted {len(pages)} pages.")
            return "\n\n".join(pages)
    except Exception as e:
        print(f"  pdfplumber failed ({e}), trying pypdf...")

    # Attempt 2: pypdf
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        pages = [p.extract_text().strip() for p in reader.pages if p.extract_text()]
        if pages:
            print(f"  pypdf extracted {len(pages)} pages.")
            return "\n\n".join(pages)
    except Exception as e:
        print(f"  pypdf failed ({e})")

    raise ValueError(
        "Could not extract text from this PDF. "
        "It may be image-based (scanned). Try a text-based PDF, .docx, or .txt file."
    )


def _parse_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n\n".join(parts)


def _parse_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _parse_csv(file_path: str) -> str:
    rows = []
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        headers = None
        for i, row in enumerate(reader):
            if i == 0:
                headers = row
                rows.append("Columns: " + ", ".join(row))
            else:
                if headers:
                    row_text = " | ".join(f"{h}: {v}" for h, v in zip(headers, row) if v.strip())
                else:
                    row_text = " | ".join(row)
                if row_text:
                    rows.append(row_text)
    return "\n".join(rows)


def _parse_xlsx(file_path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"[Sheet: {sheet.title}]")
        headers = None
        for i, row in enumerate(sheet.iter_rows(values_only=True)):
            row = [str(c) if c is not None else "" for c in row]
            if not any(row):
                continue
            if i == 0:
                headers = row
                parts.append("Columns: " + ", ".join(row))
            else:
                if headers:
                    row_text = " | ".join(f"{h}: {v}" for h, v in zip(headers, row) if v)
                else:
                    row_text = " | ".join(v for v in row if v)
                if row_text:
                    parts.append(row_text)
    return "\n".join(parts)


def _parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _parse_html(file_path: str) -> str:
    from bs4 import BeautifulSoup
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    # Remove scripts and styles
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _parse_json(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        data = json.load(f)

    def flatten(obj, prefix=""):
        lines = []
        if isinstance(obj, dict):
            for k, v in obj.items():
                lines.extend(flatten(v, f"{prefix}{k}: " if not prefix else f"{prefix} > {k}: "))
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                lines.extend(flatten(item, f"{prefix}[{i}] "))
        else:
            lines.append(f"{prefix}{obj}")
        return lines

    return "\n".join(flatten(data))
